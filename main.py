import pandas as pd
from pptx import Presentation

# -------------------------
# 1. Ler planilha
# -------------------------
df = pd.read_excel("dados.xlsx")

# -------------------------
# 2. Fazer análise
# -------------------------
# faturamento = df["valor"].sum()

media_ti = df[df["setor"] == "TI"]["salário"].mean()
media_rh = df[df["setor"] == "RH"]["salário"].mean()
media_juridico = df[df["setor"] == "jurídico"]["salário"].mean()

max_ti = df.loc[df.loc[df["setor"] == "TI"]["salário"].idxmax()]
max_juridico = df.loc[df.loc[df["setor"] == "jurídico"]["salário"].idxmax()]
max_rh = df.loc[df.loc[df["setor"] == "RH"]["salário"].idxmax()]

# Dados que vamos substituir
dados = {
    "{{media_ti}}": f"R$ {media_ti:.2f}",
    "{{media_rh}}": f"R$ {media_rh:.2f}",
    "{{media_juridico}}": f"R$ {media_juridico:.2f}",

    "{{max_ti}}": f"{max_ti['nome']} (R$ {max_ti['salário']:.2f})",
    "{{max_rh}}": f"{max_rh['nome']} (R$ {max_rh['salário']:.2f})",
    "{{max_juridico}}": f"{max_juridico['nome']} (R$ {max_juridico['salário']:.2f})",
}

# -------------------------
# 3. Abrir template pptx
# -------------------------
prs = Presentation("template.pptx")

# -------------------------
# 4. Substituir textos
# -------------------------
for slide in prs.slides:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue

        for paragraph in shape.text_frame.paragraphs:
            for chave, valor in dados.items():
                if chave in paragraph.text:
                    paragraph.text = paragraph.text.replace(chave, valor)

# -------------------------
# 5. Salvar novo relatório
# -------------------------
prs.save("relatorio_final.pptx")

print("Relatório gerado!")
